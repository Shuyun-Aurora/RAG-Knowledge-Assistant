import os
import PyPDF2
import docx
from pptx import Presentation
from pymongo import MongoClient
import gridfs
from bson import ObjectId
from config.settings import settings
from typing import Optional, Dict, Any


class DocumentRepository:
    def __init__(self):
        self.client = MongoClient(settings.MONGODB_URI)
        self.db = self.client["rag"]
        self.fs = gridfs.GridFS(self.db)

    def save_file(self, file_content: bytes, filename: str, metadata: Dict[str, Any] = None) -> str:
        """
        直接保存文件到MongoDB GridFS
        :param file_content: 文件内容
        :param filename: 文件名
        :param metadata: 元数据
        :return: 文件ID
        """
        file_id = self.fs.put(file_content, filename=filename, metadata=metadata or {})
        return str(file_id)

    def get_documents_by_course(self, course_name: str, page: int, size: int):
        query = {"metadata.course": course_name}
        total = self.db.fs.files.count_documents(query)
        skip = (page - 1) * size
        cursor = self.db.fs.files.find(query).skip(skip).limit(size).sort("uploadDate", -1)
        return {"items": list(cursor), "total": total}

    def get_file_stream(self, file_id: str):
        oid = ObjectId(file_id)
        return self.fs.get(oid)

    def get_file(self, file_id: str) -> Optional[Dict[str, Any]]:
        """
        直接从MongoDB GridFS获取文件
        :param file_id: 文件ID
        :return: 包含文件内容和元数据的字典
        """
        try:
            grid_out = self.fs.get(ObjectId(file_id))
            return {
                "content": grid_out.read(),
                "filename": grid_out.filename,
                "metadata": grid_out.metadata
            }
        except Exception:
            return None

    def delete_file(self, file_id: str) -> bool:
        """
        直接从MongoDB GridFS删除文件
        :param file_id: 文件ID
        :return: 是否删除成功
        """
        try:
            self.fs.delete(ObjectId(file_id))
            return True
        except Exception:
            return False

    def extract_text(self, file_id: str) -> str:
        """
        从数据库中获取文件并提取文本
        :param file_id: 文件ID
        :return: 提取的文本内容
        """
        file_info = self.get_file(file_id)
        if not file_info:
            raise ValueError(f"File not found with ID: {file_id}")
            
        file_bytes = file_info["content"]
        filename = file_info["filename"].lower()
        
        if filename.endswith(".pdf"):
            import io
            with io.BytesIO(file_bytes) as f:
                return "\n".join([p.extract_text() for p in PyPDF2.PdfReader(f).pages])
        elif filename.endswith(".docx"):
            import io
            with io.BytesIO(file_bytes) as f:
                doc = docx.Document(f)
            return "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith(".pptx"):
            import io
            with io.BytesIO(file_bytes) as f:
                prs = Presentation(f)
            return "\n".join([s.text for slide in prs.slides for s in slide.shapes if hasattr(s, "text")])
        else:
            raise ValueError("Unsupported file type")

    def get_all_files(self) -> list:
        """
        获取所有文件的信息
        :return: 文件信息列表
        """
        files = []
        for grid_out in self.fs.find():
            files.append({
                "_id": str(grid_out._id),
                "filename": grid_out.filename,
                "metadata": grid_out.metadata,
                "length": grid_out.length,
                "upload_date": grid_out.upload_date
            })
        return files

    def close(self):
        """关闭数据库连接"""
        self.client.close()
